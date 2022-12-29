import csv
import os
import re
import sys
import copy

import pptx
from pptx import Presentation
from pptx.util import Cm, Inches, Pt


class CreateLabel(object):
    def __init__(self, pattern_pptx, address_data):
        self.address_data = address_data
        self.prs = pptx.Presentation(pattern_pptx)
        self.slide = self.prs.slides[0]

    def create_label(self):
        self.print_name()
        self.print_address_no()
        self.print_address1()
        self.print_address2()
        return self.prs

    def print_name(self):
        name = self.create_name_string(self.address_data)
        self.replace_text('名前', name)

    def print_address_no(self):
        address_no = self.address_data['address_no']
        self.replace_text('123', address_no[:3])
        self.replace_text('4567', address_no[4:])

    def print_address1(self):
        address1 = self.address_data['address1']
        self.replace_text('住所1', address1)

    def print_address2(self):
        address2 = self.address_data['address2']
        self.replace_text('住所2', address2)

    def create_name_string(self, address_data):
        name = address_data['last_name']
        first_name_list = address_data['first_name_list']
        honorific_list = address_data['honorific_list']
        for idx, first_name in enumerate(first_name_list):
            if idx != 0:
                name += "\n" + ('　' * len(address_data['last_name']))
            name += " " + first_name + " " + honorific_list[idx]
        return name

    def replace_text(self, before, after):
        for shp in self.slide.shapes:
            if shp.has_text_frame:
                if before in shp.text:
                    new_text = re.sub(before, after, shp.text)
                    self.replace_paragraph_text_retaining_initial_formatting(
                        shp.text_frame.paragraphs[0], new_text)

    def replace_paragraph_text_retaining_initial_formatting(self, paragraph, new_text):
        p = paragraph._p
        for idx, run in enumerate(paragraph.runs):
            if idx > 0:
                p.remove(run._r)

        paragraph.runs[0].text = new_text


def read_csv(csv_file_path):
    reader = csv.reader(open(csv_file_path))
    address_list = []

    for line in reader:
        if line[0].encode().decode("utf-8") == "同上":
            address_list[-1]["first_name_list"].append(
                line[1].encode().decode("utf-8"))
            address_list[-1]["honorific_list"].append(
                line[2].encode().decode("utf-8"))
            continue

        tmp_address_data = {
            "last_name": line[0].encode().decode("utf-8"),
            "first_name_list": [line[1].encode().decode("utf-8")],
            "honorific_list": [line[2].encode().decode("utf-8")],
            "address_no": line[3].encode().decode("utf-8"),
            "address1": line[4].encode().decode("utf-8"),
            "address2": line[5].encode().decode("utf-8"),
        }

        address_list.append(tmp_address_data)
    return address_list


def shape_up_address(input_str: str) -> str:
    number_comverted = convert_number_to_kanji(input_str)
    result = convert_alphabet_half_width_to_full_width(number_comverted)
    return result


def convert_number_to_kanji(input_str: str) -> str:
    numbers = re.findall(r'\d+', input_str)
    kanji_numbers = [convert_to_kanji(int(n)) for n in numbers]
    return re.sub(r'\d+', lambda x: kanji_numbers.pop(0), input_str)


def convert_alphabet_half_width_to_full_width(input_str: str) -> str:
    output_str = ""
    for ch in input_str:
        if ch >= 'A' and ch <= 'Z':
            code_point = ord(ch)
            ch_full_width = chr(code_point + 65248)
            output_str += ch_full_width
        else:
            output_str += ch
    return output_str


def convert_to_kanji(num: int) -> str:
    kanji_numbers = ['〇', '一', '二', '三', '四', '五', '六', '七', '八', '九', '十']
    if num < 11:
        return kanji_numbers[num]
    elif num < 20:
        result = '十' + kanji_numbers[num % 10]
        return result
    elif num < 100:
        result = kanji_numbers[num // 10] + '十'
        if num % 10 != 0:
            result += kanji_numbers[num % 10]
        return result
    else:
        result = ""
        for n in str(num):
            result += kanji_numbers[int(n)]
        return result


def pptx_copy_slide(pres: pptx.Presentation, source: pptx.slide.Slide):
    dest = pres.slides.add_slide(source.slide_layout)
    for shape in dest.shapes:
        shape.element.getparent().remove(shape.element)

    for shape in source.shapes:
        new_shape = copy.deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(new_shape, 'p:extLst')

    for rel in source.part.rels:
        target = rel._target

        if "notesSlide" in rel.reltype:
            continue

        if 'chart' in rel.reltype:
            # https://github.com/scanny/python-pptx/issues/132#issuecomment-414001942
            partname = target.package.next_partname(
                pptx.parts.chart.ChartPart.partname_template)
            xlsx_blob = target.chart_workbook.xlsx_part.blob
            target = pptx.parts.chart.ChartPart(
                partname=partname,
                content_type=target.content_type,
                element=copy.deepcopy(target._element),
                package=target.package)
            target.chart_workbook.xlsx_part = pptx.parts.chart.EmbeddedXlsxPart.new(
                blob=xlsx_blob,
                package=target.package)

        if rel.is_external:
            dest.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
        else:
            dest.part.rels.get_or_add(rel.reltype, rel._target)

    return dest


def do_main():
    if len(sys.argv) < 3:
        print("")
        print("usage: print_label.py <pattern_pptx> <src_csv>")
        print("")
        sys.exit(1)

    pattern_pptx = sys.argv[1]
    src_csv_file = sys.argv[2]
    file_path_and_name = os.path.splitext(src_csv_file)[0]

    address_list = read_csv(src_csv_file)
    prs = pptx.Presentation(pattern_pptx)
    merged_prs = Presentation()
    merged_prs.slide_height = prs.slide_height
    merged_prs.slide_width = prs.slide_width

    for address_data in address_list:
        label_prs = CreateLabel(pattern_pptx, address_data)
        label = label_prs.create_label()
        pptx_copy_slide(merged_prs, label.slides[0])

    merged_prs.save(f"{file_path_and_name}.pptx")


if __name__ == u"__main__":
    do_main()
